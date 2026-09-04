"""
Parses uploaded learning material (PDF, DOCX, PPTX, TXT) into plain text,
then splits it into overlapping chunks suitable for retrieval.
"""
from __future__ import annotations
import io
import re
from dataclasses import dataclass
from typing import List

from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation


@dataclass
class Chunk:
    id: int
    text: str
    source: str


def extract_text(file_bytes: bytes, filename: str) -> str:
    """Extract raw text from an uploaded file based on its extension."""
    name = filename.lower()
    if name.endswith(".pdf"):
        return _extract_pdf(file_bytes)
    if name.endswith(".docx"):
        return _extract_docx(file_bytes)
    if name.endswith(".pptx"):
        return _extract_pptx(file_bytes)
    if name.endswith(".txt") or name.endswith(".md"):
        return file_bytes.decode("utf-8", errors="ignore")
    raise ValueError(f"Unsupported file type: {filename}")


def _extract_pdf(file_bytes: bytes) -> str:
    reader = PdfReader(io.BytesIO(file_bytes))
    pages = [page.extract_text() or "" for page in reader.pages]
    return "\n".join(pages)


def _extract_docx(file_bytes: bytes) -> str:
    doc = DocxDocument(io.BytesIO(file_bytes))
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            parts.append(" | ".join(c.text for c in row.cells))
    return "\n".join(parts)


def _extract_pptx(file_bytes: bytes) -> str:
    prs = Presentation(io.BytesIO(file_bytes))
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    parts.append(text)
    return "\n".join(parts)


def clean_text(text: str) -> str:
    text = re.sub(r"\s+", " ", text)
    return text.strip()


def chunk_text(text: str, chunk_size: int = 800, overlap: int = 150, source: str = "material") -> List[Chunk]:
    """
    Splits text into overlapping word-based chunks. Word-based (not
    char-based) chunking keeps chunks semantically coherent enough for
    retrieval without needing a sentence tokenizer/model download.
    """
    words = text.split()
    if not words:
        return []

    chunks: List[Chunk] = []
    start = 0
    idx = 0
    while start < len(words):
        end = min(start + chunk_size, len(words))
        chunk_words = words[start:end]
        chunks.append(Chunk(id=idx, text=" ".join(chunk_words), source=source))
        idx += 1
        if end == len(words):
            break
        start = end - overlap
    return chunks


def process_upload(file_bytes: bytes, filename: str) -> List[Chunk]:
    """Full pipeline: extract -> clean -> chunk."""
    raw = extract_text(file_bytes, filename)
    cleaned = clean_text(raw)
    return chunk_text(cleaned, source=filename)
